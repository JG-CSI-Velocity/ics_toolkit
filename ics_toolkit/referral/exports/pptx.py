"""Referral intelligence PPTX report -- reuses analysis PPTX infrastructure."""

from __future__ import annotations

import logging
from pathlib import Path

from ics_toolkit.analysis.analyses.base import AnalysisResult
from ics_toolkit.analysis.exports.pptx import (
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    _add_chart_slide,
    _add_section_divider,
    _add_table_slide,
    _build_analysis_lookup,
    _get_layout,
)
from ics_toolkit.settings import ReferralSettings

try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # type: ignore[assignment,misc]

logger = logging.getLogger(__name__)

REFERRAL_SECTION_MAP = {
    "Referral Intelligence Overview": ["Overview KPIs"],
    "Referrer Influence": [
        "Top Referrers",
        "Emerging Referrers",
        "Dormant High-Value Referrers",
        "One-time vs Repeat Referrers",
    ],
    "Staff & Branch": [
        "Staff Multipliers",
        "Branch Influence Density",
    ],
    "Code Health": ["Code Health Report"],
}

_MAPPED_NAMES: set[str] = set()
for _names in REFERRAL_SECTION_MAP.values():
    _MAPPED_NAMES.update(_names)


def write_referral_pptx(
    settings: ReferralSettings,
    analyses: list[AnalysisResult],
    output_path: Path | None = None,
    chart_pngs: dict[str, bytes] | None = None,
) -> Path:
    """Build and save a PPTX presentation for referral intelligence.

    Reuses slide builders from the analysis export module.
    """
    if output_path is None:
        output_path = settings.output_dir / "Referral_Intelligence.pptx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    lookup = _build_analysis_lookup(analyses)
    pngs = chart_pngs or {}

    # Create presentation
    template = settings.pptx_template
    if template and Path(template).exists():
        prs = Presentation(str(template))
        logger.info("Loaded template: %s", template)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        if template:
            logger.warning("Template not found: %s (using blank)", template)

    # Title slide
    layout = _get_layout(prs, preferred=1, fallback=0)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Referral Intelligence Report"
    subtitle = settings.client_name or (f"Client {settings.client_id or 'unknown'}")
    for ph_idx in [1, 13, 14]:
        try:
            slide.placeholders[ph_idx].text = subtitle
            break
        except (KeyError, IndexError):
            continue

    # Walk section map
    for section_name, analysis_names in REFERRAL_SECTION_MAP.items():
        section_analyses = [(name, lookup[name]) for name in analysis_names if name in lookup]
        if not section_analyses:
            continue

        _add_section_divider(prs, section_name)

        for name, analysis in section_analyses:
            _add_table_slide(prs, analysis)
            if name in pngs:
                _add_chart_slide(prs, analysis.title, pngs[name])

    # Catch unmapped analyses
    unmapped = [(a.name, a) for a in analyses if a.error is None and a.name not in _MAPPED_NAMES]
    if unmapped:
        _add_section_divider(prs, "Additional Analyses")
        for name, analysis in unmapped:
            _add_table_slide(prs, analysis)
            if name in pngs:
                _add_chart_slide(prs, analysis.title, pngs[name])

    prs.save(str(output_path))
    logger.info(
        "Referral PPTX report saved: %s (%d slides)",
        output_path,
        len(prs.slides),
    )
    return output_path
