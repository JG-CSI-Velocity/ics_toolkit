"""PowerPoint deck builder for ICS analysis report generation.

Provides:
    - DeckConfig: Pydantic model for all deck/layout settings
    - SlideContent: Dataclass defining a single slide's content
    - DeckBuilder: Class that assembles SlideContent objects into a .pptx file
    - make_figure: Helper to create consistently-sized matplotlib figures
    - apply_matplotlib_defaults: Global chart styling for matplotlib
    - setup_slide_helpers: Convenience initializer for Jupyter notebook workflows
    - inspect_template: Utility to list available layouts in a template
    - create_title_slide / create_summary_slide: Convenience factories

Dependencies:
    pip install python-pptx matplotlib
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from itertools import cycle
from pathlib import Path
from typing import Callable, Optional

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Layout positioning presets
# ---------------------------------------------------------------------------


class SinglePositioning(BaseModel):
    """Position preset for single-chart slides."""

    left: float = 0.5
    top: float = 2.5
    width: float = 6.0


class MultiPositioning(BaseModel):
    """Position preset for side-by-side chart slides."""

    top: float = 2.5
    left_pos: float = 2.3
    right_pos: float = 7.0
    width: float = 4.5


class PositioningConfig(BaseModel):
    """Layout-specific positioning overrides (inches)."""

    single_default: SinglePositioning = SinglePositioning()
    stacked_right: SinglePositioning = SinglePositioning(
        left=5.5,
        top=2.0,
        width=6.0,
    )
    split_standard: MultiPositioning = MultiPositioning(
        top=2.5,
        left_pos=2.3,
        right_pos=7.0,
        width=4.5,
    )
    split_spaced: MultiPositioning = MultiPositioning(
        top=2.5,
        left_pos=0.5,
        right_pos=7.5,
        width=4.5,
    )


# ---------------------------------------------------------------------------
# Deck configuration
# ---------------------------------------------------------------------------


class DeckConfig(BaseModel):
    """Complete configuration for deck building.

    All measurements are in inches unless stated otherwise.
    """

    # -- Chart export settings ------------------------------------------------
    dpi: int = 150
    fig_facecolor: str = "white"

    # -- Figure sizes (width, height in inches) --------------------------------
    fig_single: tuple[float, float] = (10.0, 6.0)
    fig_double: tuple[float, float] = (6.0, 4.0)
    fig_with_kpi: tuple[float, float] = (8.0, 6.0)
    fig_wide: tuple[float, float] = (12.0, 5.0)
    fig_square: tuple[float, float] = (7.0, 7.0)

    # -- Template layout indices (CSI template) --------------------------------
    #   0: Cover / Intro - Slide 1
    #   1: Cover / Intro - Slide 2        <- Title slide (text on left)
    #   2: Divider - Slide 2              <- Section divider
    #   3: Divider - Slide 3
    #   4: Analysis - Slide 1             <- Single chart
    #   5: Analysis - Slide 2             <- Single chart
    #   6: Analysis - Slide 3 - Split     <- Side-by-side
    #   7: Analysis - Slide 4 - Split     <- Side-by-side
    #   8: Analysis - Slide 5 - Thirds    <- Three charts
    #   9: Analysis - Slide 6 - Main      <- Single chart
    #  10: Analysis - Slide 7 - Stacked   <- Chart on right 2/3
    #  11: Analysis - Slide 8 - Blank     <- Flexible
    #  12: Analysis - Slide 11 - Header1
    #  13: Analysis - Slide 11 - Header2

    layout_title: int = 1
    layout_title_alt: int = 0
    layout_section: int = 2
    layout_chart: list[int] = [4, 5, 9, 11]
    layout_split: list[int] = [6, 7]
    layout_stacked: int = 10
    layout_thirds: int = 8

    # -- Layout-specific positioning -------------------------------------------
    positioning: PositioningConfig = PositioningConfig()

    # -- KPI text formatting (point sizes) -------------------------------------
    kpi_value_size: int = 28
    kpi_label_size: int = 12


# Module-level default config; callers can mutate or replace.
DECK_CONFIG = DeckConfig()


# ---------------------------------------------------------------------------
# Slide content definition
# ---------------------------------------------------------------------------


@dataclass
class SlideContent:
    """Container for all information needed to build a single slide.

    Attributes
    ----------
    slide_type:
        Determines which builder method creates this slide.
        Options: 'title', 'section', 'screenshot', 'screenshot_kpi',
                 'multi_screenshot', 'summary'.
    title:
        Text displayed in the slide's title area.
    images:
        File paths to chart/screenshot PNG images.
    kpis:
        Key-value pairs for KPI callouts. For title slides use
        ``{"subtitle": "January 2025"}``.
    bullets:
        Text items for summary slides (up to 9, displayed in 3x3 grid).
    layout_index:
        Index of slide layout from the template to use.
    """

    slide_type: str
    title: str
    images: Optional[list[str]] = None
    kpis: Optional[dict[str, str]] = None
    bullets: Optional[list[str]] = None
    layout_index: int = 5


# ---------------------------------------------------------------------------
# DeckBuilder class
# ---------------------------------------------------------------------------


class DeckBuilder:
    """Assembles SlideContent objects into a PowerPoint presentation.

    Parameters
    ----------
    template_path:
        Path to a .pptx template file.  If ``None``, a blank presentation
        is created instead.
    config:
        Deck configuration.  If ``None``, uses the module-level
        ``DECK_CONFIG`` default.
    """

    # -- Single screenshot defaults (inches) ----------------------------------
    SINGLE_IMG_LEFT = Inches(0.5)
    SINGLE_IMG_TOP = Inches(2.5)
    SINGLE_IMG_WIDTH = Inches(6)

    SINGLE_IMG_RIGHT_LEFT = Inches(5.5)
    SINGLE_IMG_RIGHT_TOP = Inches(2.0)
    SINGLE_IMG_RIGHT_WIDTH = Inches(6)

    # -- Multi screenshot - spaced layout -------------------------------------
    MULTI_IMG_TOP = Inches(2.5)
    MULTI_IMG_WIDTH = Inches(4.5)
    MULTI_IMG_LEFT_POS = Inches(0.5)
    MULTI_IMG_RIGHT_POS = Inches(7.5)

    # -- Multi screenshot - standard layout -----------------------------------
    MULTI_IMG_STD_TOP = Inches(2.5)
    MULTI_IMG_STD_LEFT_POS = Inches(2.3)
    MULTI_IMG_STD_RIGHT_POS = Inches(7.0)
    MULTI_IMG_STD_WIDTH = Inches(4.5)

    # -- Screenshot with KPIs ------------------------------------------------
    KPI_IMG_LEFT = Inches(0.3)
    KPI_IMG_TOP = Inches(1.5)
    KPI_IMG_WIDTH = Inches(6)

    KPI_TEXT_LEFT = Inches(6.8)
    KPI_TEXT_TOP_START = Inches(1.8)
    KPI_TEXT_WIDTH = Inches(2.5)
    KPI_VALUE_HEIGHT = Inches(0.5)
    KPI_LABEL_HEIGHT = Inches(0.3)
    KPI_SPACING = Inches(1.0)

    # -- Summary slide (3x3 bullet grid) --------------------------------------
    SUMMARY_COL_POSITIONS = [Inches(0.5), Inches(3.5), Inches(6.5)]
    SUMMARY_ROW_START = Inches(1.8)
    SUMMARY_ROW_SPACING = Inches(1.2)
    SUMMARY_BOX_WIDTH = Inches(2.8)
    SUMMARY_BOX_HEIGHT = Inches(1.0)

    def __init__(
        self,
        template_path: Path | str | None = None,
        config: DeckConfig | None = None,
    ) -> None:
        if template_path is not None:
            self.template_path: Path | None = Path(template_path)
        else:
            self.template_path = None
        self.config = config or DECK_CONFIG
        self.prs: Presentation | None = None

    # -- positioning helpers --------------------------------------------------

    def _get_single_positioning(self, layout_index: int) -> tuple:
        """Return (left, top, width) in EMU for a single-chart layout."""
        if layout_index == 4:
            return (Inches(0.95), Inches(2.11), Inches(5.5))
        if layout_index == 5:
            return (Inches(0.95), Inches(2.11), Inches(5.5))
        if layout_index == 9:
            return (Inches(0.80), Inches(2.12), Inches(6.0))
        if layout_index == 10:
            return (Inches(5.12), Inches(1.75), Inches(7.0))
        if layout_index == 11:
            return (Inches(2.0), Inches(2.0), Inches(8.0))
        # Default fallback
        return (Inches(0.95), Inches(2.11), Inches(5.5))

    def _get_multi_positioning(self, layout_index: int) -> tuple:
        """Return (top, left_pos, right_pos, width) in EMU for multi-chart layouts."""
        if layout_index == 6:
            return (
                Inches(2.12),
                Inches(0.80),
                Inches(6.78),
                Inches(5.5),
            )
        if layout_index == 7:
            return (
                Inches(2.10),
                Inches(0.80),
                Inches(6.78),
                Inches(5.5),
            )
        # Default
        return (
            Inches(2.12),
            Inches(0.80),
            Inches(6.78),
            Inches(5.5),
        )

    # -- build ----------------------------------------------------------------

    def build(self, slides: list[SlideContent], output_path: Path | str) -> Path:
        """Build complete PowerPoint deck from slide definitions.

        Parameters
        ----------
        slides:
            Ordered list of ``SlideContent`` objects to create.
        output_path:
            Where to save the generated .pptx file.

        Returns
        -------
        Path to the saved presentation.
        """
        output_path = Path(output_path)

        if self.template_path is not None:
            self.prs = Presentation(str(self.template_path))
            logger.info("Opened template: %s", self.template_path)
        else:
            self.prs = Presentation()
            logger.info("Created blank presentation (no template)")

        for slide_content in slides:
            self._add_slide(slide_content)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info("Deck saved: %s", output_path)

        return output_path

    # -- slide dispatch -------------------------------------------------------

    _BUILDER_REGISTRY: dict[str, str] = {
        "title": "_build_title_slide",
        "section": "_build_section_slide",
        "screenshot": "_build_screenshot_slide",
        "screenshot_kpi": "_build_screenshot_kpi_slide",
        "multi_screenshot": "_build_multi_screenshot_slide",
        "summary": "_build_summary_slide",
    }

    def _add_slide(self, content: SlideContent) -> None:
        """Create a slide and dispatch to the appropriate builder method."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[content.layout_index])

        method_name = self._BUILDER_REGISTRY.get(content.slide_type)
        if method_name is None:
            raise ValueError(
                f"Unknown slide_type: '{content.slide_type}'. "
                f"Valid types: {list(self._BUILDER_REGISTRY.keys())}"
            )

        builder = getattr(self, method_name)
        builder(slide, content)

    # -- title slide ----------------------------------------------------------

    def _build_title_slide(self, slide, content: SlideContent) -> None:
        """Build title slide with main title and optional subtitle.

        Layout 1 creates a custom text box; layout 0 uses placeholders.
        """
        if content.layout_index == 1:
            self._build_title_layout_1(slide, content)
            return

        # Layout 0 and others: use placeholders
        title_lines = content.title.split("\n") if "\n" in content.title else [content.title]

        if slide.shapes.title:
            slide.shapes.title.text = title_lines[0]

        additional_text: list[str] = title_lines[1:]
        if content.kpis and "subtitle" in content.kpis:
            additional_text.append(content.kpis["subtitle"])

        text_placeholders = [26, 29, 30, 27, 28, 31]
        for i, text in enumerate(additional_text):
            if i >= len(text_placeholders):
                break
            try:
                slide.placeholders[text_placeholders[i]].text = text
            except (KeyError, IndexError):
                try:
                    slide.placeholders[1].text = text
                except (KeyError, IndexError):
                    pass

    def _build_title_layout_1(self, slide, content: SlideContent) -> None:
        """Build layout-1 title slide with a custom text box."""
        title_lines = content.title.split("\n") if "\n" in content.title else [content.title]
        subtitle = content.kpis.get("subtitle", "") if content.kpis else ""

        full_text = title_lines[0]
        if len(title_lines) > 1:
            full_text += f"\n{title_lines[1]}"
        if len(title_lines) > 2:
            full_text += f"\n{title_lines[2]}"
        elif subtitle:
            separator = "\u2500" * 20
            full_text += f"\n{separator}\n{subtitle}"

        text_box = slide.shapes.add_textbox(
            Inches(1.0),
            Inches(3.0),
            Inches(6.0),
            Inches(2.0),
        )
        tf = text_box.text_frame
        tf.word_wrap = True

        lines = full_text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.LEFT
            if i == 0:
                p.font.size = Pt(32)
                p.font.bold = True
            else:
                p.font.size = Pt(18)

    # -- section slide --------------------------------------------------------

    def _build_section_slide(self, slide, content: SlideContent) -> None:
        """Build section divider slide with title and optional subtitle."""
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        if slide.shapes.title:
            slide.shapes.title.text = title_text

        if subtitle_text:
            for ph_idx in [44, 13, 1, 14]:
                try:
                    slide.placeholders[ph_idx].text = subtitle_text
                    break
                except (KeyError, IndexError):
                    continue

    # -- screenshot slide -----------------------------------------------------

    def _build_screenshot_slide(self, slide, content: SlideContent) -> None:
        """Build slide with title, optional subtitle, and a single image."""
        self._set_title_and_subtitle(slide, content)

        left, top, width = self._get_single_positioning(content.layout_index)

        if content.images and Path(content.images[0]).exists():
            slide.shapes.add_picture(
                content.images[0],
                left,
                top,
                width=width,
            )
        elif content.images:
            logger.warning("Image not found, skipping: %s", content.images[0])

    # -- screenshot_kpi slide -------------------------------------------------

    def _build_screenshot_kpi_slide(self, slide, content: SlideContent) -> None:
        """Build slide with image on left and KPI callouts on right."""
        self._set_title_and_subtitle(slide, content)

        left, top, width = self._get_single_positioning(content.layout_index)

        if content.images and Path(content.images[0]).exists():
            slide.shapes.add_picture(
                content.images[0],
                left,
                top,
                width=width,
            )
        elif content.images:
            logger.warning("Image not found, skipping: %s", content.images[0])

        if not content.kpis:
            return

        kpi_placeholder_pairs = [
            (26, 19),  # First KPI  (label, value)
            (27, 28),  # Second KPI (label, value)
        ]

        kpi_items = [(k, v) for k, v in content.kpis.items() if k != "subtitle"]

        for i, (label, value) in enumerate(kpi_items):
            if i >= len(kpi_placeholder_pairs):
                break

            label_idx, value_idx = kpi_placeholder_pairs[i]

            try:
                slide.placeholders[label_idx].text = label
            except (KeyError, IndexError):
                pass

            try:
                slide.placeholders[value_idx].text = str(value)
            except (KeyError, IndexError):
                pass

    # -- multi_screenshot slide -----------------------------------------------

    def _build_multi_screenshot_slide(self, slide, content: SlideContent) -> None:
        """Build slide with two images side by side."""
        self._set_title_and_subtitle(slide, content)

        top, left_pos, right_pos, width = self._get_multi_positioning(
            content.layout_index,
        )

        positions = [
            (left_pos, top, width),
            (right_pos, top, width),
        ]

        if not content.images:
            return

        for i, img_path in enumerate(content.images[:2]):
            if Path(img_path).exists():
                left, img_top, img_width = positions[i]
                slide.shapes.add_picture(img_path, left, img_top, width=img_width)
            else:
                logger.warning("Image not found, skipping: %s", img_path)

    # -- summary slide --------------------------------------------------------

    def _build_summary_slide(self, slide, content: SlideContent) -> None:
        """Build summary slide with bullets in a 3x3 grid."""
        if slide.shapes.title:
            slide.shapes.title.text = content.title

        if not content.bullets:
            return

        max_bullets = 9
        for i, bullet in enumerate(content.bullets[:max_bullets]):
            col = i % 3
            row = i // 3

            left = self.SUMMARY_COL_POSITIONS[col]
            top = self.SUMMARY_ROW_START + (row * self.SUMMARY_ROW_SPACING)

            txbox = slide.shapes.add_textbox(
                left,
                top,
                self.SUMMARY_BOX_WIDTH,
                self.SUMMARY_BOX_HEIGHT,
            )
            tf = txbox.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = bullet
            tf.paragraphs[0].font.size = Pt(14)

    # -- shared helpers -------------------------------------------------------

    @staticmethod
    def _set_title_and_subtitle(slide, content: SlideContent) -> None:
        """Parse 'Title\\nSubtitle' and set placeholders."""
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        if slide.shapes.title:
            slide.shapes.title.text = title_text

        if subtitle_text:
            try:
                slide.placeholders[13].text = subtitle_text
            except (KeyError, IndexError):
                pass


# ---------------------------------------------------------------------------
# Figure creation helper (matplotlib)
# ---------------------------------------------------------------------------


def make_figure(
    fig_type: str = "single",
    config: DeckConfig | None = None,
) -> tuple:
    """Create a matplotlib figure with standardized sizing.

    Parameters
    ----------
    fig_type:
        One of 'single', 'double', 'kpi', 'wide', 'square'.
    config:
        Deck configuration for sizes.  Falls back to ``DECK_CONFIG``.

    Returns
    -------
    (Figure, Axes) tuple.
    """
    import matplotlib.pyplot as plt

    cfg = config or DECK_CONFIG

    size_map = {
        "single": cfg.fig_single,
        "double": cfg.fig_double,
        "kpi": cfg.fig_with_kpi,
        "wide": cfg.fig_wide,
        "square": cfg.fig_square,
    }

    figsize = size_map.get(fig_type, cfg.fig_single)
    return plt.subplots(figsize=figsize)


# ---------------------------------------------------------------------------
# Matplotlib defaults
# ---------------------------------------------------------------------------


def apply_matplotlib_defaults() -> None:
    """Apply consistent matplotlib styling for slide-ready charts."""
    import matplotlib.pyplot as plt

    plt.rcParams.update(
        {
            # Font sizes
            "font.size": 12,
            "axes.titlesize": 14,
            "axes.labelsize": 12,
            "xtick.labelsize": 10,
            "ytick.labelsize": 10,
            "legend.fontsize": 10,
            "figure.titlesize": 16,
            # Visual style
            "axes.spines.top": False,
            "axes.spines.right": False,
            "axes.grid": False,
            # Colors
            "figure.facecolor": "white",
            "axes.facecolor": "white",
            "savefig.facecolor": "white",
            # Export
            "savefig.bbox": "tight",
        }
    )


# ---------------------------------------------------------------------------
# Jupyter notebook helpers
# ---------------------------------------------------------------------------


def setup_slide_helpers(
    chart_dir: Path | str,
    config: DeckConfig | None = None,
) -> tuple[list, Callable, Callable, Callable]:
    """Initialize slide list and helper functions for Jupyter workflows.

    Parameters
    ----------
    chart_dir:
        Directory where chart images will be saved.
    config:
        Deck configuration.  Falls back to ``DECK_CONFIG``.

    Returns
    -------
    (slides, add_chart_slide, add_section, add_multi_chart_slide)
    """
    import matplotlib.pyplot as plt

    cfg = config or DECK_CONFIG

    chart_dir = Path(chart_dir)
    chart_dir.mkdir(parents=True, exist_ok=True)

    slides: list[SlideContent] = []

    chart_layouts = cfg.layout_chart
    layout_cycler = cycle(chart_layouts)

    dpi = cfg.dpi
    facecolor = cfg.fig_facecolor

    # -----------------------------------------------------------------

    def add_chart_slide(
        fig,
        filename: str,
        title: str,
        slide_type: str = "screenshot",
        kpis: dict | None = None,
        layout_index: int | None = None,
    ) -> None:
        """Save matplotlib figure and add a slide to the list."""
        path = chart_dir / filename
        fig.savefig(path, dpi=dpi, bbox_inches="tight", facecolor=facecolor)
        plt.close(fig)

        actual_layout = layout_index if layout_index is not None else next(layout_cycler)

        slides.append(
            SlideContent(
                slide_type=slide_type,
                title=title,
                images=[str(path)],
                kpis=kpis,
                layout_index=actual_layout,
            )
        )

    # -----------------------------------------------------------------

    def add_section(title: str, layout_index: int | None = None) -> None:
        """Add a section divider slide."""
        actual_layout = layout_index if layout_index is not None else cfg.layout_section
        slides.append(SlideContent("section", title, layout_index=actual_layout))

    # -----------------------------------------------------------------

    def add_multi_chart_slide(
        fig1,
        fig2,
        filename1: str,
        filename2: str,
        title: str,
        layout_index: int | None = None,
    ) -> None:
        """Save two figures and add as a side-by-side slide."""
        path1 = chart_dir / filename1
        path2 = chart_dir / filename2

        fig1.savefig(path1, dpi=dpi, bbox_inches="tight", facecolor=facecolor)
        fig2.savefig(path2, dpi=dpi, bbox_inches="tight", facecolor=facecolor)

        plt.close(fig1)
        plt.close(fig2)

        actual_layout = layout_index if layout_index is not None else next(layout_cycler)

        slides.append(
            SlideContent(
                slide_type="multi_screenshot",
                title=title,
                images=[str(path1), str(path2)],
                layout_index=actual_layout,
            )
        )

    return slides, add_chart_slide, add_section, add_multi_chart_slide


# ---------------------------------------------------------------------------
# Utility functions
# ---------------------------------------------------------------------------


def inspect_template(template_path: Path | str) -> None:
    """Print all available slide layouts in a PowerPoint template.

    Useful for finding the correct ``layout_index`` values.
    """
    prs = Presentation(str(template_path))
    logger.info("Layouts in %s:", template_path)
    for i, layout in enumerate(prs.slide_layouts):
        logger.info("  %d: %s", i, layout.name)

    # Also enumerate placeholders per layout for debugging
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = [(ph.placeholder_format.idx, ph.name) for ph in layout.placeholders]
        if placeholders:
            logger.debug("  Layout %d placeholders: %s", i, placeholders)


def create_title_slide(
    client_name: str,
    period: str,
    config: DeckConfig | None = None,
    layout_index: int | None = None,
) -> SlideContent:
    """Create a configured title slide.

    Parameters
    ----------
    client_name:
        Client name for main title.
    period:
        Reporting period for subtitle (e.g., "January 2025").
    config:
        Deck configuration.  Falls back to ``DECK_CONFIG``.
    layout_index:
        Template layout override.
    """
    cfg = config or DECK_CONFIG
    actual_layout = layout_index if layout_index is not None else cfg.layout_title
    return SlideContent(
        slide_type="title",
        title=f"{client_name} Monthly Report",
        kpis={"subtitle": period},
        layout_index=actual_layout,
    )


def create_summary_slide(
    title: str,
    col1_bullets: list[str],
    col2_bullets: list[str],
    col3_bullets: list[str],
    config: DeckConfig | None = None,
    layout_index: int | None = None,
) -> SlideContent:
    """Create a 3-column summary slide.

    Bullets are interleaved row-by-row for the 3x3 grid layout:
    ``[col1[0], col2[0], col3[0], col1[1], col2[1], col3[1], ...]``
    """
    cfg = config or DECK_CONFIG
    bullets: list[str] = []
    for i in range(3):
        bullets.append(col1_bullets[i] if i < len(col1_bullets) else "")
        bullets.append(col2_bullets[i] if i < len(col2_bullets) else "")
        bullets.append(col3_bullets[i] if i < len(col3_bullets) else "")

    if layout_index is not None:
        actual_layout = layout_index
    else:
        chart_layouts = cfg.layout_chart
        actual_layout = chart_layouts[0] if chart_layouts else 5

    return SlideContent(
        slide_type="summary",
        title=title,
        bullets=bullets,
        layout_index=actual_layout,
    )
