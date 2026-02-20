"""PPTX report generator with data tables and chart images."""

import logging
from datetime import datetime
from io import BytesIO
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ics_toolkit.analysis.analyses.base import AnalysisResult
from ics_toolkit.settings import AnalysisSettings as Settings

logger = logging.getLogger(__name__)

# -- Colors ----------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x36, 0x5D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
ZEBRA_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
TOTAL_BG = RGBColor(0xE0, 0xE0, 0xE0)

# -- Table constraints -----------------------------------------------------
MAX_TABLE_ROWS = 12
MAX_TABLE_COLS = 10

# -- Slide positioning (widescreen 13.33" x 7.5") -------------------------
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

TITLE_LEFT = Inches(0.5)
TITLE_TOP = Inches(0.25)
TITLE_WIDTH = Inches(12.3)
TITLE_HEIGHT = Inches(0.55)

TABLE_LEFT = Inches(0.5)
TABLE_TOP = Inches(1.1)
TABLE_WIDTH = Inches(12.3)

CHART_LEFT = Inches(1.5)
CHART_TOP = Inches(1.1)
CHART_WIDTH = Inches(10.0)
CHART_MAX_HEIGHT = Inches(5.5)

# Maps analysis name -> section grouping for slide organization.
SECTION_MAP = {
    "Executive Summary": [
        "Executive Summary",
    ],
    "Summary": [
        "Total ICS Accounts",
        "Open ICS Accounts",
        "ICS by Stat Code",
        "Product Code Distribution",
        "Debit Distribution",
        "Debit x Prod Code",
        "Debit x Branch",
        "ICS Penetration by Branch",
    ],
    "Portfolio Health": [
        "Engagement Decay",
        "Net Portfolio Growth",
        "Spend Concentration",
        "Closure by Source",
        "Closure by Branch",
        "Closure by Account Age",
        "Net Growth by Source",
        "Closure Rate Trend",
    ],
    "Source Analysis": [
        "Source Distribution",
        "Source x Stat Code",
        "Source x Prod Code",
        "Source x Branch",
        "Account Type",
        "Source by Year",
        "Source Acquisition Mix",
    ],
    "DM Source Deep-Dive": [
        "DM Overview",
        "DM by Branch",
        "DM by Debit Status",
        "DM by Product",
        "DM by Year Opened",
        "DM Activity Summary",
        "DM Activity by Branch",
        "DM Monthly Trends",
    ],
    "REF Source Deep-Dive": [
        "REF Overview",
        "REF by Branch",
        "REF by Debit Status",
        "REF by Product",
        "REF by Year Opened",
        "REF Activity Summary",
        "REF Activity by Branch",
        "REF Monthly Trends",
    ],
    "Demographics": [
        "Age Comparison",
        "Closures",
        "Open vs Close",
        "Balance Tiers",
        "Stat Open Close",
        "Age vs Balance",
        "Balance Tier Detail",
        "Age Distribution",
        "Balance Trajectory",
    ],
    "Activity Analysis": [
        "Activity Summary",
        "Activity by Debit+Source",
        "Activity by Balance",
        "Activity by Branch",
        "Monthly Trends",
        "Activity by Source Comparison",
        "Monthly Interchange Trend",
        "Business vs Personal",
    ],
    "Cohort Analysis": [
        "Cohort Activation",
        "Cohort Heatmap",
        "Cohort Milestones",
        "Activation Summary",
        "Growth Patterns",
        "Activation Personas",
        "Branch Activation",
    ],
    "Persona Deep-Dive": [
        "Persona Overview",
        "Persona Swipe Contribution",
        "Persona by Branch",
        "Persona by Source",
        "Persona Revenue Impact",
        "Persona by Balance Tier",
        "Persona Velocity",
        "Persona Cohort Trend",
    ],
    "Performance": [
        "Days to First Use",
        "Branch Performance Index",
        "Product Code Performance",
    ],
    "Strategic Insights": [
        "Activation Funnel",
        "Revenue Impact",
        "Revenue by Branch",
        "Revenue by Source",
        "Dormant High-Balance",
    ],
}

# All analysis names referenced in SECTION_MAP (for detecting un-mapped analyses).
_MAPPED_NAMES: set[str] = set()
for _names in SECTION_MAP.values():
    _MAPPED_NAMES.update(_names)

# Strings that mark a row as a "total" row.
_TOTAL_MARKERS = {"total", "grand total"}


def _build_analysis_lookup(analyses: list[AnalysisResult]) -> dict[str, AnalysisResult]:
    """Build a name -> result lookup for successful analyses."""
    return {a.name: a for a in analyses if a.error is None}


# =========================================================================
# Public API
# =========================================================================


def write_chart_catalog(
    settings: Settings,
    analyses: list[AnalysisResult],
    chart_pngs: dict[str, bytes],
    output_path: Path | None = None,
) -> Path:
    """Build a chart catalog PPTX: one slide per chart with metadata labels.

    Each slide shows the chart index, name, rendered image, section,
    chart builder function name, and source .py file path.
    """
    import inspect

    from ics_toolkit.analysis.charts import CHART_REGISTRY

    if output_path is None:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = settings.output_dir / f"Chart_Catalog_{timestamp}.pptx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Build chart name -> metadata mapping via introspection
    pkg_root = Path(__file__).resolve().parent.parent.parent
    chart_meta: dict[str, dict[str, str]] = {}
    for name, func in CHART_REGISTRY.items():
        try:
            src = Path(inspect.getfile(func)).resolve()
            rel = src.relative_to(pkg_root)
            source_str = f"ics_toolkit/{rel}"
        except (ValueError, TypeError):
            source_str = "unknown"
        chart_meta[name] = {"function": func.__name__, "source": source_str}

    # Build name -> section lookup
    section_lookup: dict[str, str] = {}
    for section_name, names in SECTION_MAP.items():
        for n in names:
            section_lookup[n] = section_name

    # Blank widescreen presentation (no template needed for catalog)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "ICS Chart Catalog")
    subtitle = (
        f"{settings.client_name or 'Client ' + (settings.client_id or 'unknown')}"
        f"  --  {len(chart_pngs)} charts"
    )
    txbox = slide.shapes.add_textbox(TITLE_LEFT, Inches(1.2), TITLE_WIDTH, Inches(1.0))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT

    # One slide per chart
    for idx, (name, png_bytes) in enumerate(chart_pngs.items(), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title: "#idx  name"
        _add_slide_title(slide, f"#{idx}  {name}")

        # Chart image (reduced height to leave room for metadata)
        slide.shapes.add_picture(
            BytesIO(png_bytes),
            CHART_LEFT,
            CHART_TOP,
            width=CHART_WIDTH,
            height=Inches(4.5),
        )

        # Metadata text box
        meta = chart_meta.get(name, {})
        section = section_lookup.get(name, "Unmapped")
        func_name = meta.get("function", "N/A")
        source = meta.get("source", "N/A")

        meta_lines = [
            f"Section: {section}",
            f"Function: {func_name}()",
            f"Source: {source}",
        ]

        txbox = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(5.8),
            Inches(12.3),
            Inches(1.2),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(meta_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(str(output_path))
    logger.info("Chart catalog saved: %s (%d slides)", output_path, len(prs.slides))
    return output_path


def write_pptx_report(
    settings: Settings,
    analyses: list[AnalysisResult],
    output_path: Path | None = None,
    chart_pngs: dict[str, bytes] | None = None,
) -> Path:
    """Build and save a PPTX presentation with tables and charts.

    Every successful analysis gets a data-table slide.
    Analyses with chart PNGs get an additional chart slide.
    """
    if output_path is None:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = settings.output_dir / f"ICS_Analysis_{timestamp}.pptx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    lookup = _build_analysis_lookup(analyses)
    pngs = chart_pngs or {}

    prs = _create_presentation(settings)

    # Title slide
    _add_title_slide(prs, settings)

    # Walk SECTION_MAP in order
    for section_name, analysis_names in SECTION_MAP.items():
        section_analyses = [(name, lookup[name]) for name in analysis_names if name in lookup]
        if not section_analyses:
            continue

        _add_section_divider(prs, section_name)

        for name, analysis in section_analyses:
            # Table slide for every analysis (shows the actual data)
            _add_table_slide(prs, analysis)

            # Chart slide if PNG available
            if name in pngs:
                _add_chart_slide(prs, analysis.title, pngs[name])

    # Catch any analyses NOT in SECTION_MAP (shouldn't lose data)
    unmapped = [(a.name, a) for a in analyses if a.error is None and a.name not in _MAPPED_NAMES]
    if unmapped:
        _add_section_divider(prs, "Additional Analyses")
        for name, analysis in unmapped:
            _add_table_slide(prs, analysis)
            if name in pngs:
                _add_chart_slide(prs, analysis.title, pngs[name])

    prs.save(str(output_path))
    logger.info("PPTX report saved: %s (%d slides)", output_path, len(prs.slides))
    return output_path


# =========================================================================
# Presentation setup
# =========================================================================


def _create_presentation(settings: Settings) -> Presentation:
    """Load CSI template or create blank widescreen presentation."""
    template = settings.pptx_template
    if template and Path(template).exists():
        prs = Presentation(str(template))
        logger.info("Loaded template: %s", template)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        if template:
            logger.warning("Template not found: %s (using blank)", template)
    return prs


def _get_layout(prs: Presentation, preferred: int, fallback: int = 6):
    """Get slide layout by index with fallback."""
    for idx in [preferred, fallback, 6, 5, 0]:
        try:
            return prs.slide_layouts[idx]
        except IndexError:
            continue
    return prs.slide_layouts[0]


# =========================================================================
# Slide builders
# =========================================================================


def _add_title_slide(prs: Presentation, settings: Settings) -> None:
    """Add branded title slide."""
    layout = _get_layout(prs, preferred=1, fallback=0)
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = "ICS Accounts Analysis"

    subtitle = settings.client_name or f"Client {settings.client_id}"
    for ph_idx in [1, 13, 14]:
        try:
            slide.placeholders[ph_idx].text = subtitle
            break
        except (KeyError, IndexError):
            continue


def _add_section_divider(prs: Presentation, title: str) -> None:
    """Add section divider slide."""
    layout = _get_layout(prs, preferred=2, fallback=5)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title


def _add_table_slide(prs: Presentation, analysis: AnalysisResult) -> None:
    """Add slide with analysis title and formatted data table."""
    df = analysis.df
    if df.empty:
        return

    layout = _get_layout(prs, preferred=11, fallback=6)
    slide = prs.slides.add_slide(layout)

    # Title
    _add_slide_title(slide, analysis.title)

    # Prepare data for table (truncate if needed)
    show_df, truncated = _prepare_table_df(df)
    cols = list(show_df.columns[:MAX_TABLE_COLS])
    col_truncated = len(df.columns) > MAX_TABLE_COLS

    nrows = len(show_df) + 1  # +1 for header
    ncols = len(cols)

    # Calculate table height: ~0.45" per row (Pt 20 font)
    row_height = 0.45
    table_height = Inches(nrows * row_height)

    tbl_shape = slide.shapes.add_table(
        nrows,
        ncols,
        TABLE_LEFT,
        TABLE_TOP,
        TABLE_WIDTH,
        table_height,
    )
    table = tbl_shape.table

    # Distribute column widths
    col_width = int(TABLE_WIDTH / ncols)
    for j in range(ncols):
        table.columns[j].width = col_width

    # Header row
    for j, col_name in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        _style_header_cell(cell)

    # Data rows
    for i, (_, row) in enumerate(show_df.iterrows()):
        is_total = _is_total_row(row)
        is_odd = i % 2 == 1

        for j, col_name in enumerate(cols):
            cell = table.cell(i + 1, j)
            cell.text = _format_cell_value(row[col_name], col_name)
            _style_data_cell(cell, is_total=is_total, is_odd=is_odd)

    # Truncation note
    if truncated or col_truncated:
        parts = []
        if truncated:
            parts.append(f"Showing {len(show_df)} of {len(df)} rows")
        if col_truncated:
            parts.append(f"{len(cols)} of {len(df.columns)} columns")
        note = " | ".join(parts) + " -- see Excel report for full data"
        note_top = TABLE_TOP + table_height + Inches(0.15)
        _add_footnote(slide, note, note_top)


def _add_chart_slide(prs: Presentation, title: str, png_bytes: bytes) -> None:
    """Add slide with chart image."""
    layout = _get_layout(prs, preferred=11, fallback=6)
    slide = prs.slides.add_slide(layout)

    _add_slide_title(slide, title)

    slide.shapes.add_picture(
        BytesIO(png_bytes),
        CHART_LEFT,
        CHART_TOP,
        width=CHART_WIDTH,
        height=CHART_MAX_HEIGHT,
    )


# =========================================================================
# Shared slide helpers
# =========================================================================


def _add_slide_title(slide, text: str) -> None:
    """Add a title text box to the slide."""
    txbox = slide.shapes.add_textbox(
        TITLE_LEFT,
        TITLE_TOP,
        TITLE_WIDTH,
        TITLE_HEIGHT,
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY


def _add_footnote(slide, text: str, top) -> None:
    """Add a small footnote text box."""
    txbox = slide.shapes.add_textbox(
        TABLE_LEFT,
        top,
        TABLE_WIDTH,
        Inches(0.3),
    )
    p = txbox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


# =========================================================================
# Table helpers
# =========================================================================


def _prepare_table_df(df: pd.DataFrame) -> tuple[pd.DataFrame, bool]:
    """Truncate DataFrame for slide display, preserving total rows.

    Returns (display_df, was_truncated).
    """
    if len(df) <= MAX_TABLE_ROWS:
        return df, False

    # Preserve total/grand total rows at bottom
    total_mask = df.iloc[:, 0].astype(str).str.strip().str.lower().isin(_TOTAL_MARKERS)
    total_rows = df[total_mask]

    if len(total_rows) > 0:
        n_total = len(total_rows)
        data_rows = df[~total_mask].head(MAX_TABLE_ROWS - n_total)
        return pd.concat([data_rows, total_rows]), True

    return df.head(MAX_TABLE_ROWS), True


def _is_total_row(row: pd.Series) -> bool:
    """Detect if a row is a Total/Grand Total summary row."""
    first = str(row.iloc[0]).strip().lower()
    return first in _TOTAL_MARKERS


def _format_cell_value(val, col_name: str) -> str:
    """Format a cell value for slide display."""
    if pd.isna(val):
        return ""

    col_lower = col_name.lower()

    # Percentage columns
    if "%" in col_name or "rate" in col_lower or "pct" in col_lower:
        if isinstance(val, (int, float)):
            return f"{val:.1f}%"

    # Currency columns
    if any(kw in col_lower for kw in ("balance", "spend", "revenue", "interchange")):
        if isinstance(val, (int, float)):
            if abs(val) >= 1000:
                return f"${val:,.0f}"
            return f"${val:,.2f}"

    # Integer counts
    if isinstance(val, (int,)):
        return f"{val:,}"

    # Floats
    if isinstance(val, float):
        if val == int(val) and abs(val) < 1e9:
            return f"{int(val):,}"
        return f"{val:,.1f}"

    return str(val)


def _style_header_cell(cell) -> None:
    """Style a table header cell: navy background, white bold text."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY

    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER


def _style_data_cell(
    cell,
    is_total: bool = False,
    is_odd: bool = False,
) -> None:
    """Style a data cell with alternating shading and total highlighting."""
    if is_total:
        cell.fill.solid()
        cell.fill.fore_color.rgb = TOTAL_BG
    elif is_odd:
        cell.fill.solid()
        cell.fill.fore_color.rgb = ZEBRA_GRAY
    else:
        cell.fill.background()

    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = DARK_TEXT
        paragraph.alignment = PP_ALIGN.CENTER
        if is_total:
            paragraph.font.bold = True
