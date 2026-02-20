"""Tests for exports/pptx.py -- PPTX report generation."""

import pandas as pd
import pytest
from pptx import Presentation

from ics_toolkit.analysis.analyses.base import AnalysisResult
from ics_toolkit.analysis.exports.pptx import (
    SECTION_MAP,
    _build_analysis_lookup,
    _format_cell_value,
    _is_total_row,
    _prepare_table_df,
    write_pptx_report,
)


@pytest.fixture
def mock_analyses():
    """Create mock analysis results matching SECTION_MAP names."""
    return [
        AnalysisResult(
            name="Total ICS Accounts",
            title="ICS Accounts - Total Dataset",
            df=pd.DataFrame({"Branch": ["Main", "West", "Total"], "Count": [60, 40, 100]}),
            sheet_name="01_Total",
        ),
        AnalysisResult(
            name="Open ICS Accounts",
            title="ICS Accounts - Open Accounts",
            df=pd.DataFrame({"Branch": ["Main", "Total"], "Count": [50, 50]}),
            sheet_name="02_Open",
        ),
        AnalysisResult(
            name="Bad Analysis",
            title="Failed",
            df=pd.DataFrame(),
            error="Something broke",
        ),
    ]


class TestBuildAnalysisLookup:
    def test_builds_lookup(self, mock_analyses):
        lookup = _build_analysis_lookup(mock_analyses)
        assert "Total ICS Accounts" in lookup
        assert "Open ICS Accounts" in lookup

    def test_excludes_errored(self, mock_analyses):
        lookup = _build_analysis_lookup(mock_analyses)
        assert "Bad Analysis" not in lookup

    def test_empty_list(self):
        assert _build_analysis_lookup([]) == {}


class TestFormatCellValue:
    def test_nan_returns_dash(self):
        assert _format_cell_value(float("nan"), "Count") == "-"

    def test_percentage_column(self):
        assert _format_cell_value(45.678, "Debit %") == "45.7%"
        assert _format_cell_value(100.0, "Activation Rate") == "100.0%"

    def test_currency_column(self):
        assert _format_cell_value(1234.56, "Avg Balance") == "$1,235"
        assert _format_cell_value(50.25, "Avg Spend") == "$50.25"

    def test_integer(self):
        assert _format_cell_value(1234, "Count") == "1,234"

    def test_float_whole_number(self):
        assert _format_cell_value(100.0, "Count") == "100"

    def test_string(self):
        assert _format_cell_value("Main", "Branch") == "Main"


class TestIsTotalRow:
    def test_total(self):
        row = pd.Series(["Total", 100])
        assert _is_total_row(row)

    def test_grand_total(self):
        row = pd.Series(["Grand Total", 200])
        assert _is_total_row(row)

    def test_not_total(self):
        row = pd.Series(["Main", 50])
        assert not _is_total_row(row)


class TestPrepareTableDf:
    def test_small_df_unchanged(self):
        df = pd.DataFrame({"A": range(5)})
        result, truncated = _prepare_table_df(df)
        assert len(result) == 5
        assert not truncated

    def test_large_df_truncated(self):
        df = pd.DataFrame({"A": [f"Row {i}" for i in range(50)], "B": range(50)})
        result, truncated = _prepare_table_df(df)
        assert len(result) <= 20
        assert truncated

    def test_preserves_total_row(self):
        rows = [f"Branch {i}" for i in range(30)] + ["Grand Total"]
        df = pd.DataFrame({"Branch": rows, "Count": range(31)})
        result, truncated = _prepare_table_df(df)
        assert truncated
        last_val = str(result.iloc[-1, 0]).strip().lower()
        assert last_val == "grand total"


class TestWritePptxReport:
    def test_creates_file(self, sample_settings, mock_analyses, tmp_path):
        path = tmp_path / "test.pptx"
        result = write_pptx_report(sample_settings, mock_analyses, output_path=path)
        assert result.exists()

    def test_auto_generates_path(self, sample_settings, mock_analyses):
        sample_settings.output_dir.mkdir(parents=True, exist_ok=True)
        result = write_pptx_report(sample_settings, mock_analyses)
        assert result.exists()
        assert "ICS_Analysis" in result.name

    def test_no_analyses(self, sample_settings, tmp_path):
        path = tmp_path / "no_analyses.pptx"
        result = write_pptx_report(sample_settings, [], output_path=path)
        assert result.exists()

    def test_has_title_slide(self, sample_settings, mock_analyses, tmp_path):
        path = tmp_path / "test.pptx"
        write_pptx_report(sample_settings, mock_analyses, output_path=path)
        prs = Presentation(str(path))
        assert len(prs.slides) >= 1

    def test_includes_table_slides(self, sample_settings, tmp_path):
        """Each successful non-empty analysis should produce a table slide."""
        analyses = [
            AnalysisResult(
                name="Total ICS Accounts",
                title="Total ICS Accounts",
                df=pd.DataFrame({"Branch": ["Main"], "Count": [100]}),
            ),
            AnalysisResult(
                name="ICS by Stat Code",
                title="ICS by Stat Code",
                df=pd.DataFrame({"Stat": ["O", "C"], "Count": [80, 20]}),
            ),
        ]
        path = tmp_path / "tables.pptx"
        write_pptx_report(sample_settings, analyses, output_path=path)
        prs = Presentation(str(path))
        # Title + section divider + 2 table slides = at least 4
        assert len(prs.slides) >= 4

    def test_chart_pngs_add_slides(self, sample_settings, tmp_path):
        """Analyses with chart PNGs get an extra chart slide."""
        analyses = [
            AnalysisResult(
                name="Total ICS Accounts",
                title="Total ICS Accounts",
                df=pd.DataFrame({"Branch": ["Main"], "Count": [100]}),
            ),
        ]
        # Create a minimal valid PNG (1x1 pixel)
        import struct
        import zlib

        def _make_tiny_png() -> bytes:
            raw = b"\x00\xff\x00\x00"
            compressed = zlib.compress(raw)
            ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)

            def chunk(tag, data):
                c = tag + data
                return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

            return (
                b"\x89PNG\r\n\x1a\n"
                + chunk(b"IHDR", ihdr)
                + chunk(b"IDAT", compressed)
                + chunk(b"IEND", b"")
            )

        pngs = {"Total ICS Accounts": _make_tiny_png()}
        path = tmp_path / "with_charts.pptx"
        write_pptx_report(
            sample_settings, analyses, output_path=path, chart_pngs=pngs,
        )
        prs = Presentation(str(path))
        # Title + section + table + chart = at least 4
        assert len(prs.slides) >= 4

    def test_unmapped_analyses_included(self, sample_settings, tmp_path):
        """Analyses not in SECTION_MAP still appear in the deck."""
        analyses = [
            AnalysisResult(
                name="Custom Analysis XYZ",
                title="Custom Analysis XYZ",
                df=pd.DataFrame({"Metric": ["A"], "Value": [42]}),
            ),
        ]
        path = tmp_path / "unmapped.pptx"
        write_pptx_report(sample_settings, analyses, output_path=path)
        prs = Presentation(str(path))
        # Title + "Additional Analyses" divider + table = 3
        assert len(prs.slides) >= 3


class TestSectionMap:
    def test_has_expected_sections(self):
        expected = {
            "Summary",
            "Source Analysis",
            "DM Source Deep-Dive",
            "Demographics",
            "Activity Analysis",
            "Cohort Analysis",
            "Performance",
            "Portfolio Health",
            "Strategic Insights",
            "Executive Summary",
        }
        assert set(SECTION_MAP.keys()) == expected

    def test_all_lists_nonempty(self):
        for section, names in SECTION_MAP.items():
            assert len(names) > 0, f"Section {section} is empty"

    def test_activity_summary_name_matches_registry(self):
        """Verify Activity Summary uses the correct registry name."""
        activity = SECTION_MAP["Activity Analysis"]
        assert "Activity Summary" in activity
        assert "L12M Activity Summary" not in activity
